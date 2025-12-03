from pptx import Presentation
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR  
from pptx.util import Pt

def slayd_generator(mavzu, ism, rejalars, matnlar):
    input_file = "pptx/shablon_2.pptx"
    output_file = "edited_shablon.pptx"
    presentation = Presentation(input_file)

    # Birinchi slaydda mavzu va tayyorlovchini yangilash
    first_slide = presentation.slides[0]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if "Mavzu:MavzuMavzu" == paragraph.text:
                    paragraph.text = f"Mavzu: {mavzu}"
                elif "Tayyorladi:" == paragraph.text:
                    paragraph.text = f"Tayyorladi: {ism}"

    # Ikkinchi slaydda rejalarni yangilash
    reja_slide = presentation.slides[1]
    for shape in reja_slide.shapes:
        if shape.has_text_frame and "RejalarRejalarRejalar" in shape.text_frame.text:
            text_frame = shape.text_frame
            text_frame.clear()
            for reja in rejalars:
                text_frame.add_paragraph().text = reja

    # Har bir reja uchun yangi slaydlar yaratish
    for reja, matn in matnlar.items():
        slide_layout = presentation.slide_layouts[1]  # Sarlavha va matn uchun slayd
        is_first_slide = True  # Sarlavhani faqat birinchi slaydga yozish uchun flag

        while matn:
            # Matnni kesib olish (har bir bo'lak 724 ta belgidan oshmasligi kerak)
            chunk = matn[:724]
            matn = matn[724:]

            # Yangi slayd qo'shish
            new_slide = presentation.slides.add_slide(slide_layout)

            # Agar bu birinchi slayd bo'lsa sarlavha qo'shish, keyingi slaydlarda olib tashlash
            if is_first_slide:
                title = new_slide.shapes.title
                title.text = reja
                title.text_frame.paragraphs[0].font.size = Pt(27)
                title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                is_first_slide = False  # Keyingi slaydlarda sarlavha yozilmaydi
            else:
                # Sarlavha shaklini butunlay o'chirish
                for shape in new_slide.shapes:
                    if shape == new_slide.shapes.title:
                        sp = shape
                        sp._element.getparent().remove(sp._element)

            # Matn joylashtirish va markazga hizlash
            content = new_slide.placeholders[1]
            content.text = chunk
            text_frame = content.text_frame
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in text_frame.paragraphs:
                paragraph.font.size = Pt(23)
                paragraph.alignment = PP_ALIGN.CENTER

    # Taqdimotni saqlash
    presentation.save(output_file)
    return output_file


# Rejalar va matnlar ro'yxati
rejalar = [
    "1. Birinchi Reja",
    "2. Ikkinchi Reja",
    "3. Uchinchi Reja",
    "Xulosa"

]
matnlar = {
    "1. Birinchi Reja": "llakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljsk",
    "2. Ikkinchi Reja": "llakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljsk",
    "3. Uchinchi Reja": "llakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljsk",
    "Xulosa": "llakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljskllakjdlakjsdaljsk"

    
}

output_file = slayd_generator(mavzu="Yangi taqdimot", ism="Amirjon Karimov", rejalars=rejalar, matnlar=matnlar)
print(f"Tahrir qilingan taqdimot '{output_file}' nomi bilan saqlandi.")
